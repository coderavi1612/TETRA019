# pyrefly: ignore [missing-import]
from pptx import Presentation
import os
import datetime
# pyrefly: ignore [missing-import]
import pptx

from app.schemas.parsed_document import ParsedDocument, ContentBlock, BlockSource, SectionInfo, DocumentMetadata, ParserInfo, DocumentStatistics
from app.parsers.base import BaseParser

class PPTParser(BaseParser):
    def parse(self, file_path: str, company_id: str, document_type: str) -> ParsedDocument:
        doc_name = os.path.basename(file_path)
        file_size = os.path.getsize(file_path)
        from app.core import get_utc_now_iso
        created_at = get_utc_now_iso()
        
        # 1. Run basic loop over slide shapes using python-pptx to get raw strings
        prs = Presentation(file_path)
        slides_count = len(prs.slides)
        
        slides_raw_data = []
        for slide_idx, slide in enumerate(prs.slides):
            slide_num = slide_idx + 1
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    t = shape.text.strip()
                    if t:
                        texts.append(t)
                elif shape.has_table:
                    for r in range(len(shape.table.rows)):
                        row_cells = []
                        for c in range(len(shape.table.columns)):
                            cell_t = shape.table.cell(r, c).text.strip()
                            if cell_t:
                                row_cells.append(cell_t)
                        if row_cells:
                            texts.append(" | ".join(row_cells))
            slides_raw_data.append({
                "slide_number": slide_num,
                "raw_text": "\n".join(texts)
            })

        # 2. Call Gemini to perform structured text clean up and formatting
        api_key = os.getenv("GEMINI_API_KEY") or os.getenv("GOOGLE_API_KEY")
        if not api_key:
            raise ValueError("GEMINI_API_KEY environment variable is not set.")
            
        from google import genai
        from google.genai import types
        from pydantic import BaseModel
        import json

        class SlideText(BaseModel):
            slide_number: int
            text: str

        class PresentationText(BaseModel):
            slides: list[SlideText]

        model_name = os.getenv("GEMINI_MODEL", "gemini-3.5-flash")
        client = genai.Client(api_key=api_key)

        prompt = (
            "You are an expert document text extractor. You are provided with raw slide-by-slide text dumped from a PowerPoint presentation. "
            "Clean up, refine, and extract all text, tables (format them as markdown tables), and layout content slide-by-slide. "
            "Do not infer, estimate, or perform calculations. Keep numbers and names exact. "
            "Return the structured presentation text content organized by slide number."
        )

        res = client.models.generate_content(
            model=model_name,
            contents=[
                json.dumps(slides_raw_data, indent=2),
                prompt
            ],
            config=types.GenerateContentConfig(
                response_mime_type="application/json",
                response_schema=PresentationText,
            )
        )

        pres_text_data = json.loads(res.text)
        pres_text = PresentationText(**pres_text_data)

        content_blocks = []
        sequence = 0
        words_count = 0

        for s in pres_text.slides:
            content_blocks.append(ContentBlock(
                id=f"{document_type}_slide_{s.slide_number:02d}_visual_block",
                sequence=sequence,
                content_type="text",
                slide=s.slide_number,
                raw_text=s.text,
                section=SectionInfo(),
                source=BlockSource(file=doc_name, slide=s.slide_number)
            ))
            sequence += 1
            words_count += len(s.text.split())

        statistics = DocumentStatistics(
            pages=0,
            slides=slides_count,
            sheets=0,
            blocks=len(content_blocks),
            tables=0,
            words=words_count
        )
        
        metadata = DocumentMetadata(
            company_id=company_id,
            file_size=file_size,
            extension=".pptx",
            mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            created_at=created_at,
            parser=ParserInfo(name="Gemini-PPTX-Extractor", version="1.0"),
            statistics=statistics
        )
        
        document_id = f"{company_id}_{document_type}"
        
        return ParsedDocument(
            document_id=document_id,
            document_name=doc_name,
            document_type=document_type,
            status="parsed",
            metadata=metadata,
            content=content_blocks,
            warnings=[],
            errors=[]
        )
